"""
Export Service
--------------
Converts a Resource's structured_output JSON into PDF, DOCX, or PPTX files.

PDF  – WeasyPrint (renders an HTML template to PDF)
DOCX – python-docx (programmatic Word document)
PPTX – python-pptx (programmatic PowerPoint, slides only)
"""

import base64
import io
import json
import re
from typing import Any

from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Colour schemes for PPTX export
# ---------------------------------------------------------------------------

PPTX_SCHEMES = {
    "ocean": {
        "name": "Ocean",
        "bg":     (13,  42,  65),
        "panel":  (22,  65,  98),
        "accent": (95,  156, 179),
        "accent2":(162, 210, 228),
        "title_c":(255, 255, 255),
        "body_c": (200, 230, 240),
        "muted":  (120, 170, 195),
    },
    "midnight": {
        "name": "Midnight",
        "bg":     (10,  10,  28),
        "panel":  (18,  18,  55),
        "accent": (110, 120, 255),
        "accent2":(180, 185, 255),
        "title_c":(255, 255, 255),
        "body_c": (200, 200, 240),
        "muted":  (130, 130, 200),
    },
    "forest": {
        "name": "Forest",
        "bg":     (10,  32,  20),
        "panel":  (18,  55,  35),
        "accent": (72,  175, 110),
        "accent2":(150, 220, 175),
        "title_c":(255, 255, 255),
        "body_c": (195, 235, 215),
        "muted":  (110, 185, 145),
    },
    "charcoal": {
        "name": "Charcoal",
        "bg":     (22,  22,  28),
        "panel":  (40,  40,  52),
        "accent": (255, 140, 50),
        "accent2":(255, 200, 140),
        "title_c":(255, 255, 255),
        "body_c": (215, 215, 225),
        "muted":  (145, 145, 165),
    },
    "crimson": {
        "name": "Crimson",
        "bg":     (32,  8,   16),
        "panel":  (58,  12,  28),
        "accent": (210, 60,  80),
        "accent2":(255, 160, 175),
        "title_c":(255, 255, 255),
        "body_c": (240, 200, 210),
        "muted":  (185, 120, 140),
    },
    "clean": {
        "name": "Clean",
        "bg":     (255, 255, 255),
        "panel":  (235, 245, 252),
        "accent": (64,  101, 123),
        "accent2":(95,  156, 179),
        "title_c":(20,  40,  60),
        "body_c": (50,  75,  95),
        "muted":  (110, 140, 160),
    },
}


# ---------------------------------------------------------------------------
# PDF export – WeasyPrint
# ---------------------------------------------------------------------------

def export_pdf(resource: Any) -> bytes:
    """Render resource to HTML then convert to PDF bytes via WeasyPrint."""
    try:
        from weasyprint import HTML, CSS
    except ImportError as exc:
        raise RuntimeError(
            "WeasyPrint is not installed or its system dependencies are missing. "
            "Install WeasyPrint and ensure libcairo / libpango are available."
        ) from exc

    data = json.loads(resource.structured_output)
    html_content = _render_print_html(resource, data)
    css = CSS(string=_print_css())
    return HTML(string=html_content).write_pdf(stylesheets=[css])


def _print_css() -> str:
    return """
    @page { margin: 2cm; size: A4; }
    body { font-family: 'Helvetica Neue', Arial, sans-serif; font-size: 11pt; color: #1f2937; line-height: 1.5; }
    h1 { font-size: 20pt; color: #1e3a5f; margin-bottom: 4pt; }
    h2 { font-size: 14pt; color: #1e40af; margin-top: 16pt; margin-bottom: 6pt; border-bottom: 1pt solid #93c5fd; padding-bottom: 3pt; }
    h3 { font-size: 12pt; color: #374151; margin-top: 10pt; margin-bottom: 4pt; }
    .meta { color: #6b7280; font-size: 9pt; margin-bottom: 12pt; }
    .section-card { border: 1pt solid #e5e7eb; border-radius: 6pt; padding: 10pt; margin-bottom: 10pt; }
    .question { margin-bottom: 8pt; }
    .marks { color: #6b7280; font-style: italic; font-size: 9pt; }
    .answer-space { border-bottom: 1pt solid #d1d5db; margin-top: 4pt; height: 12pt; }
    ul { margin: 4pt 0; padding-left: 18pt; }
    li { margin-bottom: 2pt; }
    table { width: 100%; border-collapse: collapse; margin-top: 8pt; }
    th { background: #eff6ff; text-align: left; padding: 6pt; font-size: 10pt; border: 1pt solid #bfdbfe; }
    td { padding: 6pt; border: 1pt solid #e5e7eb; font-size: 10pt; vertical-align: top; }
    .notes { background: #f9fafb; border-left: 3pt solid #93c5fd; padding: 6pt 10pt; font-size: 9pt; color: #4b5563; margin-top: 4pt; }
    """


def _render_print_html(resource: Any, data: dict) -> str:
    """Build an HTML string representing the resource for PDF rendering."""
    rtype = resource.type

    sections_html = ""
    if rtype == "lesson":
        sections_html = _lesson_html(data)
    elif rtype == "worksheet":
        sections_html = _worksheet_html(data)
    elif rtype == "scheme":
        sections_html = _scheme_html(data)
    elif rtype == "slides":
        sections_html = _slides_html(data)

    return f"""<!DOCTYPE html>
<html lang="en">
<head><meta charset="utf-8"><title>{resource.title}</title></head>
<body>
  <h1>{resource.title}</h1>
  <p class="meta">{resource.subject} &bull; {resource.key_stage} &bull; {resource.topic}</p>
  {sections_html}
</body>
</html>"""


def _lesson_html(data: dict) -> str:
    objs = "".join(f"<li>{o}</li>" for o in data.get("learning_objectives", []))
    resources = "".join(f"<li>{r}</li>" for r in data.get("resources_needed", []))
    diff = data.get("differentiation", {})

    sections = ""
    for s in data.get("sections", []):
        notes = f'<div class="notes">Teacher notes: {s["teacher_notes"]}</div>' if s.get("teacher_notes") else ""
        sections += f"""
        <div class="section-card">
          <h3>{s.get('title','')} <small style="color:#9ca3af;">({s.get('duration','')})</small></h3>
          <p>{s.get('activity','')}</p>
          {notes}
        </div>"""

    return f"""
    <h2>Learning Objectives</h2><ul>{objs}</ul>
    <h2>Lesson Sections</h2>{sections}
    <h2>Resources Needed</h2><ul>{resources}</ul>
    <h2>Differentiation</h2>
    <p><strong>Support:</strong> {diff.get('support','')}</p>
    <p><strong>Extension:</strong> {diff.get('extension','')}</p>
    <h2>Assessment</h2><p>{data.get('assessment','')}</p>"""


def _worksheet_html(data: dict) -> str:
    html = f"<p><strong>Instructions:</strong> {data.get('instructions','')}</p>"
    for section in data.get("sections", []):
        qs = ""
        for q in section.get("questions", []):
            lines = "".join(f'<div class="answer-space"></div>' for _ in range(q.get("answer_lines", 3)))
            qs += f"""
            <div class="question">
              <p><strong>Q{q['number']}.</strong> {q['question']}
                <span class="marks">({q.get('marks',1)} mark{'s' if q.get('marks',1)>1 else ''})</span>
              </p>
              {lines}
            </div>"""
        html += f"""
        <div class="section-card">
          <h3>{section.get('title','')}</h3>
          {qs}
        </div>"""
    return html


def _scheme_html(data: dict) -> str:
    html = f"<p>{data.get('overview','')}</p><h2>Weekly Overview</h2>"
    html += "<table><tr><th>Week</th><th>Topic</th><th>Objectives</th><th>Activities</th><th>Assessment</th></tr>"
    for w in data.get("weeks", []):
        objs = "<br>".join(w.get("objectives", []))
        acts = "<br>".join(w.get("key_activities", []))
        html += f"<tr><td><strong>Week {w['week']}</strong></td><td>{w.get('topic','')}</td><td>{objs}</td><td>{acts}</td><td>{w.get('assessment','')}</td></tr>"
    html += "</table>"
    return html


def _slides_html(data: dict) -> str:
    html = f"<p>{data.get('slide_count', 0)} slides</p>"
    for slide in data.get("slides", []):
        bullets = "".join(f"<li>{b}</li>" for b in slide.get("bullet_points", []))
        notes = f'<div class="notes">Speaker notes: {slide["speaker_notes"]}</div>' if slide.get("speaker_notes") else ""
        html += f"""
        <div class="section-card">
          <h3>Slide {slide['number']}: {slide.get('title','')}</h3>
          <ul>{bullets}</ul>
          {notes}
        </div>"""
    return html


# ---------------------------------------------------------------------------
# DOCX export – python-docx
# ---------------------------------------------------------------------------

def export_docx(resource: Any) -> bytes:
    data = json.loads(resource.structured_output)
    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    rtype = resource.type
    if rtype == "lesson":
        _docx_lesson(doc, resource, data)
    elif rtype == "worksheet":
        _docx_worksheet(doc, resource, data)
    elif rtype == "scheme":
        _docx_scheme(doc, resource, data)
    else:
        _docx_generic(doc, resource, data)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _add_heading(doc: Document, text: str, level: int = 1) -> None:
    doc.add_heading(text, level=level)


def _docx_lesson(doc: Document, resource: Any, data: dict) -> None:
    _add_heading(doc, data.get("title", resource.title), 1)
    doc.add_paragraph(
        f"Subject: {resource.subject}  |  Key Stage: {resource.key_stage}  "
        f"|  Duration: {data.get('duration', '60 minutes')}"
    ).runs[0].italic = True

    _add_heading(doc, "Learning Objectives", 2)
    for obj in data.get("learning_objectives", []):
        doc.add_paragraph(obj, style="List Bullet")

    _add_heading(doc, "Lesson Sections", 2)
    for s in data.get("sections", []):
        doc.add_heading(f"{s.get('title','')}  ({s.get('duration','')})", 3)
        doc.add_paragraph(s.get("activity", ""))
        if s.get("teacher_notes"):
            note_para = doc.add_paragraph(f"Teacher notes: {s['teacher_notes']}")
            note_para.runs[0].italic = True
            note_para.runs[0].font.color.rgb = RGBColor(75, 85, 99)

    _add_heading(doc, "Resources Needed", 2)
    for r in data.get("resources_needed", []):
        doc.add_paragraph(r, style="List Bullet")

    diff = data.get("differentiation", {})
    if diff:
        _add_heading(doc, "Differentiation", 2)
        doc.add_paragraph(f"Support: {diff.get('support', '')}")
        doc.add_paragraph(f"Extension: {diff.get('extension', '')}")

    _add_heading(doc, "Assessment", 2)
    doc.add_paragraph(data.get("assessment", ""))


def _docx_worksheet(doc: Document, resource: Any, data: dict) -> None:
    _add_heading(doc, data.get("title", resource.title), 1)
    doc.add_paragraph(
        f"Subject: {resource.subject}  |  Key Stage: {resource.key_stage}  |  Topic: {resource.topic}"
    ).runs[0].italic = True
    doc.add_paragraph()
    doc.add_paragraph(f"Instructions: {data.get('instructions', '')}").runs[0].bold = True

    q_num = 0
    for section in data.get("sections", []):
        _add_heading(doc, section.get("title", ""), 2)
        for q in section.get("questions", []):
            q_num += 1
            marks = q.get("marks", 1)
            para = doc.add_paragraph()
            run = para.add_run(f"Q{q['number']}. ")
            run.bold = True
            para.add_run(f"{q['question']} ({marks} mark{'s' if marks > 1 else ''})")
            for _ in range(q.get("answer_lines", 3)):
                line_para = doc.add_paragraph("_" * 80)
                line_para.paragraph_format.space_before = Pt(2)
                line_para.paragraph_format.space_after = Pt(2)
                line_para.runs[0].font.color.rgb = RGBColor(209, 213, 219)


def _docx_scheme(doc: Document, resource: Any, data: dict) -> None:
    _add_heading(doc, data.get("title", resource.title), 1)
    doc.add_paragraph(
        f"Subject: {resource.subject}  |  Key Stage: {resource.key_stage}  |  Duration: {data.get('duration','6 weeks')}"
    ).runs[0].italic = True
    doc.add_paragraph(data.get("overview", ""))

    _add_heading(doc, "Weekly Overview", 2)
    table = doc.add_table(rows=1, cols=5)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, h in enumerate(["Week", "Topic", "Objectives", "Activities", "Assessment"]):
        hdr[i].text = h
        hdr[i].paragraphs[0].runs[0].bold = True

    for w in data.get("weeks", []):
        row = table.add_row().cells
        row[0].text = str(w.get("week", ""))
        row[1].text = w.get("topic", "")
        row[2].text = "\n".join(w.get("objectives", []))
        row[3].text = "\n".join(w.get("key_activities", []))
        row[4].text = w.get("assessment", "")


def _docx_generic(doc: Document, resource: Any, data: dict) -> None:
    _add_heading(doc, data.get("title", resource.title), 1)
    doc.add_paragraph(json.dumps(data, indent=2))


# ---------------------------------------------------------------------------
# PPTX export – python-pptx  (professional themed slides)
# ---------------------------------------------------------------------------

# ── Low-level helpers ──────────────────────────────────────────────────────

def _r(t: tuple) -> PptRGB:
    return PptRGB(*t)


def _rect(slide, x: float, y: float, w: float, h: float, color: tuple):
    """Add a filled rectangle (no border). Coordinates in inches."""
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = _r(color)
    s.line.fill.background()
    return s


def _tb(slide, x: float, y: float, w: float, h: float):
    """Create a word-wrapping textbox and return its text_frame."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _run(para, text: str, size: int, color: tuple, bold=False, italic=False):
    r = para.add_run()
    r.text = text
    r.font.size = PptPt(size)
    r.font.color.rgb = _r(color)
    r.font.bold = bold
    r.font.italic = italic
    return r


def _single(slide, x, y, w, h, text, size, color, bold=False,
            align=PP_ALIGN.LEFT, italic=False):
    tf = _tb(slide, x, y, w, h)
    p = tf.paragraphs[0]
    p.alignment = align
    _run(p, text, size, color, bold, italic)


def _bullets(slide, x, y, w, h, items, size, color, lead="▸  "):
    if not items:
        return
    tf = _tb(slide, x, y, w, h)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = PptPt(9)
        p.space_after  = PptPt(2)
        _run(p, f"{lead}{item}", size, color)


def _slide_num(slide, num, color):
    _single(slide, 12.55, 7.08, 0.65, 0.32, str(num), 11, color,
            align=PP_ALIGN.RIGHT)


def _set_notes(slide, text: str):
    if text and text.strip():
        slide.notes_slide.notes_text_frame.text = text.strip()


def _embed_image(slide, x, y, w, h, data_url: str):
    """Decode a data-URL and place the image on the slide."""
    if not data_url or not data_url.startswith("data:"):
        return
    m = re.match(r"data:[^;]+;base64,(.+)", data_url, re.DOTALL)
    if not m:
        return
    try:
        raw = base64.b64decode(m.group(1))
        slide.shapes.add_picture(
            io.BytesIO(raw), Inches(x), Inches(y),
            width=Inches(w), height=Inches(h),
        )
    except Exception:
        pass  # gracefully skip broken image data


# ── Slide layout builders ──────────────────────────────────────────────────

_W, _H = 13.33, 7.5   # slide dimensions in inches


def _pptx_title_slide(slide, data: dict, sc: dict):
    """Full-height accent panel on left, large title on right."""
    _rect(slide, 0,   0,   _W,  _H,  sc["bg"])        # background
    _rect(slide, 0,   0,   4.2, _H,  sc["accent"])    # left accent panel
    _rect(slide, 4.5, 4.05, 8.5, 0.04, sc["accent2"]) # thin divider

    _single(slide, 4.5, 1.7, 8.5, 2.1,
            data.get("title", ""), 36, sc["title_c"], bold=True)

    bullets = [b for b in data.get("bullet_points", []) if b.strip()]
    if bullets:
        _bullets(slide, 4.5, 4.25, 8.5, 2.6, bullets, 19, sc["accent2"], lead="")

    _slide_num(slide, data.get("number", 1), sc["muted"])
    _set_notes(slide, data.get("speaker_notes", ""))


def _pptx_content_slide(slide, data: dict, sc: dict):
    """Top accent bar, large title, bullet list, optional embedded image."""
    _rect(slide, 0,   0,    _W,   _H,   sc["bg"])      # background
    _rect(slide, 0,   0,    _W,   0.16, sc["accent"])  # top bar
    _rect(slide, 0.5, 1.18, 12.3, 0.03, sc["accent"])  # title divider

    _single(slide, 0.5, 0.26, 12.3, 0.9,
            data.get("title", ""), 26, sc["title_c"], bold=True)

    img_url = data.get("image", "")
    content_w = 7.6 if img_url else 12.3

    bullets = [b for b in data.get("bullet_points", []) if b.strip()]
    _bullets(slide, 0.5, 1.35, content_w, 5.45, bullets, 17, sc["body_c"])

    if img_url:
        _embed_image(slide, 8.4, 1.35, 4.6, 5.45, img_url)

    _slide_num(slide, data.get("number", ""), sc["muted"])
    _set_notes(slide, data.get("speaker_notes", ""))


def _pptx_image_section_slide(slide, data: dict, sc: dict):
    """Stand-alone image slide (from legacy image-section nodes)."""
    _rect(slide, 0,   0,   _W,   _H,   sc["bg"])
    _rect(slide, 0,   0,   _W,   0.16, sc["accent"])
    caption = data.get("caption") or data.get("title") or "Image"
    _single(slide, 0.5, 0.26, 12.3, 0.9, caption, 22, sc["title_c"],
            bold=True, align=PP_ALIGN.CENTER)
    src = data.get("src") or data.get("image", "")
    if src:
        _embed_image(slide, 1.5, 1.3, 10.33, 5.5, src)


# ── Public export function ─────────────────────────────────────────────────

def export_pptx(resource: Any) -> bytes:
    data = json.loads(resource.structured_output)
    slides_data = data.get("slides", [])
    if not slides_data:
        raise ValueError(
            "No slides found. PPTX export is only available for Slide Outline resources."
        )

    sc = PPTX_SCHEMES.get(
        data.get("colour_scheme", "ocean"), PPTX_SCHEMES["ocean"]
    )

    prs = Presentation()
    prs.slide_width  = Inches(_W)
    prs.slide_height = Inches(_H)
    blank = prs.slide_layouts[6]  # completely blank

    for s in slides_data:
        slide = prs.slides.add_slide(blank)
        if s.get("type") == "image":
            _pptx_image_section_slide(slide, s, sc)
        elif s.get("content_type") == "title":
            _pptx_title_slide(slide, s, sc)
        else:
            _pptx_content_slide(slide, s, sc)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
